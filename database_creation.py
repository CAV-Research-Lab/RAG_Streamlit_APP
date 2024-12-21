import os
import PyPDF2
from docx import Document
from pptx import Presentation
import json

# __extracting contents from files__ 

def extract_text_from_pdf(pdf_path): # Extracts text from a PDF file
    """Extracts text from a PDF file."""
    text = ""
    with open(pdf_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + '\n'
    return text


def extract_text_from_docx(docx_path): # Extracts text from a Word file
    """Extracts text from a Word file."""
    text = ""
    document = Document(docx_path)
    for paragraph in document.paragraphs:
        text += paragraph.text + '\n'
    return text


def extract_text_from_pptx(pptx_path): # Extracts text from a PowerPoint file
    """Extracts text from a PowerPoint file."""
    text = ""
    presentation = Presentation(pptx_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + '\n'
    return text


def save_text_with_metadata(output_path, text, metadata): # Saves extracted text with metadata into a txt file
    """Saves extracted text with metadata into a txt file."""
    with open(output_path, 'w', encoding='utf-8') as file:
        file.write("=== Metadata ===\n")
        file.write(json.dumps(metadata, indent=4))
        file.write("\n\n=== Content ===\n")
        file.write(text)


# Main function to process files
def process_files(input_dir, output_dir): # The process_files function processes files in a specified input directory and saves the extracted text with metadata into a text file in the output directory. The extract_text_from_pdf, extract_text_from_docx, and extract_text_from_pptx functions are used to extract text from PDF, Word, and PowerPoint files, respectively. The save_text_with_metadata function saves the extracted text along with metadata (e.g., file name, file path, status) into a text file. The main function processes files in the input directory and saves the extracted text with metadata into the output directory.
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    for file_name in os.listdir(input_dir):
        file_path = os.path.join(input_dir, file_name)
        output_path = os.path.join(output_dir, f"{os.path.splitext(file_name)[0]}.txt")

        if file_name.endswith('.pdf'):
            text = extract_text_from_pdf(file_path)
        elif file_name.endswith('.docx'):
            text = extract_text_from_docx(file_path)
        elif file_name.endswith('.pptx'):
            text = extract_text_from_pptx(file_path)
        else:
            print(f"Skipping unsupported file format: {file_name}")
            continue

        # Metadata with file name included
        metadata = {
            "file_name": file_name,
            "file_path": file_path,
            "status": "Processed"
        }

        # Save the extracted text with metadata
        save_text_with_metadata(output_path, text, metadata)
        print(f"Finished processing: {file_name}")


if __name__ == "__main__":
    # Get the directory of the current script
    current_dir = os.path.dirname(__file__)
    input_directory = os.path.join(current_dir, "resources") # input directory path 
    output_directory = os.path.join(current_dir, "saved_data", "database")  # Output directory path
    process_files(input_directory, output_directory) # The process_files function processes files in a specified input directory and saves the extracted text with metadata into a text file in the output directory. The extract_text_from_pdf, extract_text_from_docx, and extract_text_from_pptx functions are used to extract text from PDF, Word, and PowerPoint files, respectively. The save_text_with_metadata function saves the extracted text along with metadata (e.g., file name, file path, status) into a text file. The main function processes files in the input directory and saves the extracted text with metadata into the output directory.